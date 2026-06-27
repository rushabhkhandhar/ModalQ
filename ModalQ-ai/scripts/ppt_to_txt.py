import io
from pptx import Presentation

def ppt_to_text(ppt_file):
    text = ""
    try:
        # Load the presentation
        prs = Presentation(ppt_file)
        
        for i, slide in enumerate(prs.slides):
            text += f"Slide {i+1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + "\n"
            text += "\n"
            
            # Print progress
            # print(f"Processed slide {i+1}/{len(prs.slides)}")
    
    except Exception as e:
        print(f"An error occurred: {str(e)}")
    
    return text

# Usage
# ppt_file = "/home/yuvraj/Coding/OverloadOblivion_Datahack/data/ch 1.pptx"
# text = ppt_to_text(ppt_file)

# # Print first 1000 characters to preview
# print(text[:10000000])

# # Save the entire text to a file
# with io.open("extracted_text.txt", "w", encoding="utf-8") as f:
#     f.write(text)

# print(f"Full text has been saved to 'extracted_text.txt'")